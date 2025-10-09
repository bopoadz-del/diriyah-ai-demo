"""Utility helpers for extracting text from common construction file types."""
from __future__ import annotations
import zipfile
from pathlib import Path
from typing import Iterable
import xml.etree.ElementTree as ET

try:  # pragma: no cover - optional PDF dependency
    import PyPDF2
except ImportError:  # pragma: no cover - handled gracefully
    PyPDF2 = None  # type: ignore[assignment]

try:  # pragma: no cover - optional DOCX dependency
    import docx
except ImportError:  # pragma: no cover - handled gracefully
    docx = None  # type: ignore[assignment]

try:  # pragma: no cover - optional XLSX dependency
    import openpyxl
except ImportError:  # pragma: no cover - handled gracefully
    openpyxl = None  # type: ignore[assignment]

try:  # pragma: no cover - optional PPTX dependency
    import pptx
except ImportError:  # pragma: no cover - handled gracefully
    pptx = None  # type: ignore[assignment]

try:  # pragma: no cover - optional archive dependency
    import rarfile
except ImportError:  # pragma: no cover - handled gracefully
    rarfile = None  # type: ignore[assignment]

def _read_text_file(path: Path) -> str:
    with path.open("r", encoding="utf-8", errors="ignore") as handle:
        return handle.read()

def _extract_from_pdf(path: Path) -> str:
    if PyPDF2 is None:
        return "PDF text extraction unavailable (PyPDF2 not installed)"
    with path.open("rb") as handle:
        reader = PyPDF2.PdfReader(handle)
        return "".join(page.extract_text() or "" for page in reader.pages)

def _extract_from_docx(path: Path) -> str:
    if docx is None:
        return "DOCX extraction unavailable (python-docx not installed)"
    document = docx.Document(path)
    return "\n".join(paragraph.text for paragraph in document.paragraphs)

def _extract_from_pptx(path: Path) -> str:
    if pptx is None:
        return "PPTX extraction unavailable (python-pptx not installed)"
    presentation = pptx.Presentation(path)
    text_segments: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_segments.append(shape.text)
    return "\n".join(text_segments)

def _extract_from_xlsx(path: Path) -> str:
    if openpyxl is None:
        return "XLSX extraction unavailable (openpyxl not installed)"
    workbook = openpyxl.load_workbook(path, data_only=True)
    rows: list[str] = []
    for sheet_name in workbook.sheetnames:
        worksheet = workbook[sheet_name]
        for row in worksheet.iter_rows(values_only=True):
            cells = [str(cell) for cell in row if cell not in (None, "") ]
            if cells:
                rows.append(" ".join(cells))
    return "\n".join(rows)

def _extract_from_archive(path: Path, *, rar: bool = False) -> str:
    text_segments: list[str] = []
    if rar:
        if rarfile is None:
            return "RAR extraction unavailable (rarfile not installed)"
        archive_cls = rarfile.RarFile
    else:
        archive_cls = zipfile.ZipFile
    with archive_cls(path) as archive:
        names: Iterable[str] = archive.namelist()
        for name in names:
            try:
                data = archive.read(name)
                text_segments.append(data.decode("utf-8"))
            except Exception:
                continue
    return "\n".join(text_segments)

def extract_file_content(filepath: str) -> str:
    """Extract text content from ``filepath``.

    Favour resilience so callers always receive a string payload—even on failure."""
    path = Path(filepath)
    suffix = path.suffix.lower()
    try:
        if suffix == ".pdf":
            return _extract_from_pdf(path)
        if suffix == ".docx":
            return _extract_from_docx(path)
        if suffix == ".pptx":
            return _extract_from_pptx(path)
        if suffix == ".xlsx":
            return _extract_from_xlsx(path)
        if suffix == ".xml":
            root = ET.parse(path).getroot()
            return ET.tostring(root, encoding="unicode")
        if suffix == ".xer":
            return _read_text_file(path)
        if suffix == ".zip":
            return _extract_from_archive(path)
        if suffix == ".rar":
            return _extract_from_archive(path, rar=True)
        return _read_text_file(path)
    except Exception as exc:  # pragma: no cover - defensive logging path
        return f"[ERROR parsing {filepath}: {exc}]"
